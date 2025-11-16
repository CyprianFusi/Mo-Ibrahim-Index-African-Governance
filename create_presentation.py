"""
PowerPoint Presentation Generator for IIAG Analysis
Creates a professional 15-20 minute presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import numpy as np
from pathlib import Path

print("Creating PowerPoint Presentation...")
print("=" * 60)

# Load data for slide content
data_path = Path('data/csv-files')
composite_scores = pd.read_csv(data_path / '2024 IIAG_Composite Scores.csv', encoding='utf-8-sig')
composite_scores = composite_scores.replace('.', np.nan)
for col in composite_scores.columns[3:]:
    composite_scores[col] = pd.to_numeric(composite_scores[col], errors='coerce')

# Regional groupings
regional_groups = {
    'North Africa': ['Algeria', 'Egypt', 'Libya', 'Morocco', 'Tunisia'],
    'West Africa': ['Benin', 'Burkina Faso', 'Cape Verde', 'Côte d\'Ivoire', 'Gambia', 'Ghana', 'Guinea',
                    'Guinea-Bissau', 'Liberia', 'Mali', 'Mauritania', 'Niger', 'Nigeria', 'Senegal',
                    'Sierra Leone', 'Togo'],
    'East Africa': ['Burundi', 'Comoros', 'Djibouti', 'Eritrea', 'Ethiopia', 'Kenya', 'Madagascar',
                    'Mauritius', 'Rwanda', 'Seychelles', 'Somalia', 'South Sudan', 'Sudan', 'Tanzania', 'Uganda'],
    'Central Africa': ['Cameroon', 'Central African Republic', 'Chad', 'Congo', 'DR Congo', 'Equatorial Guinea',
                       'Gabon', 'São Tomé and Príncipe'],
    'Southern Africa': ['Angola', 'Botswana', 'Eswatini', 'Lesotho', 'Malawi', 'Mozambique', 'Namibia',
                        'South Africa', 'Zambia', 'Zimbabwe']
}

def get_region(country):
    for region, countries in regional_groups.items():
        if country in countries:
            return region
    return 'Other'

composite_scores['Region'] = composite_scores['Country'].apply(get_region)
latest_year = composite_scores['Year'].max()
latest_data = composite_scores[composite_scores['Year'] == latest_year].copy()

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(102, 126, 234)  # Blue
SECONDARY_COLOR = RGBColor(118, 75, 162)  # Purple
ACCENT_COLOR = RGBColor(46, 204, 113)  # Green
TEXT_COLOR = RGBColor(44, 62, 80)  # Dark gray
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE

    return slide

def add_content_slide(prs, title, content_type='bullet'):
    """Add a content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR

    # Underline
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()

    return slide

def add_bullet_points(slide, bullets, start_top=1.5, start_left=0.8, width=8.5):
    """Add bullet points to a slide"""
    text_box = slide.shapes.add_textbox(Inches(start_left), Inches(start_top), Inches(width), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)

    return text_box

def add_image_slide(prs, title, image_path):
    """Add a slide with an image"""
    slide = add_content_slide(prs, title)

    # Add image
    img_path = Path(image_path)
    if img_path.exists():
        left = Inches(0.5)
        top = Inches(1.3)
        slide.shapes.add_picture(str(img_path), left, top, width=Inches(9))

    return slide

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
print("  [1/20] Creating title slide...")
add_title_slide(prs,
                "African Governance Landscape",
                "Ibrahim Index Analysis 2014-2023")

# ============================================================================
# SLIDE 2: Agenda
# ============================================================================
print("  [2/20] Creating agenda slide...")
slide = add_content_slide(prs, "Presentation Agenda")
bullets = [
    "Overview: What is the Ibrahim Index?",
    "Continental Snapshot: Current State of Governance",
    "Top Performers: Success Stories",
    "Improvement Champions: Countries Making Progress",
    "Challenges: Areas Needing Attention",
    "Regional Analysis: Geographic Patterns",
    "Category Deep Dive: What Drives Good Governance?",
    "Recommendations & Next Steps"
]
add_bullet_points(slide, bullets)

# ============================================================================
# SLIDE 3: What is IIAG?
# ============================================================================
print("  [3/20] Creating IIAG overview slide...")
slide = add_content_slide(prs, "What is the Ibrahim Index of African Governance?")
bullets = [
    "Most comprehensive assessment of African governance",
    "Covers 54 countries across the continent (2014-2023)",
    "Measures 4 key categories with 16 subcategories",
    "Scores range from 0-100 (100 = best governance)",
    "Data from 100+ indicators (World Bank, UN, AfDB, etc.)",
    "Published annually by Mo Ibrahim Foundation"
]
add_bullet_points(slide, bullets)

# ============================================================================
# SLIDE 4: Four Governance Categories
# ============================================================================
print("  [4/20] Creating categories slide...")
slide = add_content_slide(prs, "Four Pillars of African Governance")
bullets = [
    "1. Security & Rule of Law - Safety, justice, accountability, anti-corruption",
    "2. Participation, Rights & Inclusion - Democracy, rights, equality",
    "3. Foundations for Economic Opportunity - Infrastructure, business environment",
    "4. Human Development - Health, education, social protection, environment"
]
add_bullet_points(slide, bullets, start_top=1.8)

# Stats box
stats_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(1.5))
stats_frame = stats_box.text_frame
stats_text = f"2023 Continental Averages:\nSecurity & Rule of Law: 47.9  |  Participation & Rights: 48.7\nEconomic Opportunity: 48.9  |  Human Development: 51.6"
stats_frame.text = stats_text
stats_para = stats_frame.paragraphs[0]
stats_para.font.size = Pt(16)
stats_para.font.color.rgb = SECONDARY_COLOR
stats_para.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 5: Key Statistics
# ============================================================================
print("  [5/20] Creating key statistics slide...")
slide = add_content_slide(prs, "Continental Snapshot (2023)")

# Create stat boxes
stats = [
    ("54", "African\nCountries"),
    ("49.3", "Continental\nAverage"),
    ("75.3", "Highest Score\n(Seychelles)"),
    ("19.0", "Lowest Score\n(South Sudan)")
]

box_width = 2
box_height = 1.8
start_left = 0.7
start_top = 2.5
spacing = 2.2

for i, (number, label) in enumerate(stats):
    left = start_left + (i * spacing)

    # Box background
    box = slide.shapes.add_shape(1, Inches(left), Inches(start_top), Inches(box_width), Inches(box_height))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY_COLOR if i % 2 == 0 else SECONDARY_COLOR
    box.line.fill.background()

    # Number
    num_box = slide.shapes.add_textbox(Inches(left), Inches(start_top + 0.3), Inches(box_width), Inches(0.8))
    num_frame = num_box.text_frame
    num_frame.text = number
    num_para = num_frame.paragraphs[0]
    num_para.alignment = PP_ALIGN.CENTER
    num_para.font.size = Pt(44)
    num_para.font.bold = True
    num_para.font.color.rgb = WHITE

    # Label
    label_box = slide.shapes.add_textbox(Inches(left), Inches(start_top + 1.1), Inches(box_width), Inches(0.6))
    label_frame = label_box.text_frame
    label_frame.text = label
    label_para = label_frame.paragraphs[0]
    label_para.alignment = PP_ALIGN.CENTER
    label_para.font.size = Pt(14)
    label_para.font.color.rgb = WHITE

# Add context
context_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
context_frame = context_box.text_frame
context_frame.text = "56.3-point spread demonstrates vast governance diversity across Africa"
context_para = context_frame.paragraphs[0]
context_para.alignment = PP_ALIGN.CENTER
context_para.font.size = Pt(18)
context_para.font.color.rgb = TEXT_COLOR
context_para.font.italic = True

# ============================================================================
# SLIDE 6: Top Performers Visualization
# ============================================================================
print("  [6/20] Creating top performers visualization...")
add_image_slide(prs, "Top & Bottom Performers (2023)", "visualizations/02_top_bottom_countries.png")

# ============================================================================
# SLIDE 7: Top 10 Analysis
# ============================================================================
print("  [7/20] Creating top 10 analysis slide...")
slide = add_content_slide(prs, "Excellence in Governance: Top 10 Countries")

top_10 = latest_data.nlargest(10, 'OVERALL GOVERNANCE')[['Country', 'OVERALL GOVERNANCE']]

# Create table
rows = 11
cols = 2
left = Inches(2)
top = Inches(1.8)
width = Inches(6)
height = Inches(4.5)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(4)
table.columns[1].width = Inches(2)

# Header
table.cell(0, 0).text = "Country"
table.cell(0, 1).text = "Score"

for i in range(2):
    cell = table.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY_COLOR
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(18)

# Data rows
for i, (idx, row) in enumerate(top_10.iterrows(), 1):
    table.cell(i, 0).text = f"{i}. {row['Country']}"
    table.cell(i, 1).text = f"{row['OVERALL GOVERNANCE']:.1f}"

    for j in range(2):
        cell = table.cell(i, j)
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

# ============================================================================
# SLIDE 8: Top Performers Insights
# ============================================================================
print("  [8/20] Creating top performers insights...")
slide = add_content_slide(prs, "What Makes Top Performers Succeed?")
bullets = [
    "Island Nation Advantage: Top 3 are all islands (Seychelles, Mauritius, Cabo Verde)",
    "  - Manageable scale of governance",
    "  - Economic specialization (tourism, services)",
    "  - Strong institutions and political stability",
    "",
    "Mainland Leaders: South Africa (4th), Botswana (5th) demonstrate continental success",
    "",
    "Common Success Factors:",
    "  - Strong rule of law and judicial independence",
    "  - Low corruption levels",
    "  - High human development investments",
    "  - Democratic governance and civic participation"
]
add_bullet_points(slide, bullets)

# ============================================================================
# SLIDE 9: Temporal Trends
# ============================================================================
print("  [9/20] Creating temporal trends visualization...")
add_image_slide(prs, "10-Year Governance Trends (2014-2023)", "visualizations/03_temporal_trends.png")

# ============================================================================
# SLIDE 10: Improvement Champions
# ============================================================================
print("  [10/20] Creating improvement champions slide...")
slide = add_content_slide(prs, "Remarkable Improvements: Top Gainers (2014-2023)")

# Calculate changes
def calculate_change(df, country, start_year, end_year):
    start_score = df[(df['Country'] == country) & (df['Year'] == start_year)]['OVERALL GOVERNANCE'].values
    end_score = df[(df['Country'] == country) & (df['Year'] == end_year)]['OVERALL GOVERNANCE'].values
    if len(start_score) > 0 and len(end_score) > 0:
        return end_score[0] - start_score[0]
    return np.nan

countries_list = composite_scores['Country'].unique()
changes = []
for country in countries_list:
    change = calculate_change(composite_scores, country, 2014, 2023)
    if not np.isnan(change):
        changes.append({'Country': country, 'Change': change})

changes_df = pd.DataFrame(changes).sort_values('Change', ascending=False)
top_improvers = changes_df.head(8)

# Create table
rows = 9
cols = 2
left = Inches(2)
top = Inches(1.8)
width = Inches(6)
height = Inches(4.2)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(4)
table.columns[1].width = Inches(2)

# Header
table.cell(0, 0).text = "Country"
table.cell(0, 1).text = "Improvement"

for i in range(2):
    cell = table.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_COLOR
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(18)

# Data
for i, (idx, row) in enumerate(top_improvers.iterrows(), 1):
    table.cell(i, 0).text = f"{i}. {row['Country']}"
    table.cell(i, 1).text = f"+{row['Change']:.1f}"

    for j in range(2):
        cell = table.cell(i, j)
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR if j == 1 else TEXT_COLOR
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 250, 240)

# ============================================================================
# SLIDE 11: Success Stories
# ============================================================================
print("  [11/20] Creating success stories slide...")
slide = add_content_slide(prs, "Success Stories: What Drove Improvement?")
bullets = [
    "Seychelles (+10.0 points): Sustained economic diversification and institutional strengthening",
    "",
    "Gambia (+7.2 points): Democratic transition after 22-year autocracy",
    "  - 2016 election marked peaceful transfer of power",
    "  - Restoration of civil liberties and press freedom",
    "",
    "Somalia (+6.8 points): Post-conflict recovery and state-building",
    "  - Federal government formation",
    "  - Security sector improvements",
    "",
    "Key Lesson: Transformation is possible with political will and sustained reform"
]
add_bullet_points(slide, bullets)

# ============================================================================
# SLIDE 12: Governance Change - All Countries
# ============================================================================
print("  [12/20] Creating governance change visualization...")
add_image_slide(prs, "Improvers vs. Decliners: Complete Picture", "visualizations/05_governance_change_all.png")

# ============================================================================
# SLIDE 13: Challenges & Decliners
# ============================================================================
print("  [13/20] Creating challenges slide...")
slide = add_content_slide(prs, "Governance Challenges: Countries in Decline")

bottom_decliners = changes_df.tail(6)

bullets = []
for idx, row in bottom_decliners.iterrows():
    bullets.append(f"{row['Country']}: {row['Change']:.1f} points")

bullets.insert(0, "Concerning trends in previously strong performers:")
bullets.append("")
bullets.append("Common factors driving decline:")
bullets.append("  - Democratic backsliding and erosion of civil liberties")
bullets.append("  - Economic pressures reducing governance capacity")
bullets.append("  - Security deterioration (especially Sahel region)")
bullets.append("  - Institutional complacency after years of strong performance")

add_bullet_points(slide, bullets)

# ============================================================================
# SLIDE 14: Regional Analysis
# ============================================================================
print("  [14/20] Creating regional analysis visualization...")
add_image_slide(prs, "Regional Patterns: Geographic Governance Divide", "visualizations/06_regional_comparison.png")

# ============================================================================
# SLIDE 15: Regional Insights
# ============================================================================
print("  [15/20] Creating regional insights slide...")
slide = add_content_slide(prs, "Regional Performance Analysis")

regional_stats = latest_data.groupby('Region')['OVERALL GOVERNANCE'].agg(['mean', 'count']).round(1)
regional_stats = regional_stats[regional_stats.index != 'Other'].sort_values('mean', ascending=False)

bullets = [
    f"Southern Africa leads at {regional_stats.loc['Southern Africa', 'mean']} average",
    "  - Benefits from SADC integration and relatively stable democracies",
    "  - Strong performers: Botswana, Namibia, South Africa",
    "",
    f"Central Africa faces challenges at {regional_stats.loc['Central Africa', 'mean']} average",
    "  - Lowest regional score - 15 points below leaders",
    "  - Persistent conflicts and weak institutions",
    "",
    f"East Africa shows highest variation (includes both Seychelles at 75.3 and South Sudan at 19.0)",
    "",
    "Key Insight: Regional cooperation matters - integrated regions perform better"
]
add_bullet_points(slide, bullets)

# ============================================================================
# SLIDE 16: Category Performance
# ============================================================================
print("  [16/20] Creating category heatmap...")
add_image_slide(prs, "Multi-Dimensional View: Category Performance", "visualizations/04_category_heatmap_top20.png")

# ============================================================================
# SLIDE 17: Category Correlations
# ============================================================================
print("  [17/20] Creating category correlation visualization...")
add_image_slide(prs, "What Drives Overall Governance?", "visualizations/07_category_correlation.png")

# ============================================================================
# SLIDE 18: Key Insights
# ============================================================================
print("  [18/20] Creating key insights slide...")
slide = add_content_slide(prs, "Key Insights & Takeaways")
bullets = [
    "1. Island nations excel - but mainland success is possible (South Africa, Botswana)",
    "",
    "2. Reform momentum works - Gambia, Somalia show dramatic improvements possible",
    "",
    "3. Vigilance required - even strong performers can decline (Tunisia, Mauritius)",
    "",
    "4. Regional disparities persist - 15-point gap between best and worst regions",
    "",
    "5. All categories matter - strong correlations show interconnected reform needed",
    "",
    "6. Human Development leads - highest average score (51.6)",
    "",
    "7. Security & Rule of Law most variable - critical foundation for progress"
]
add_bullet_points(slide, bullets)

# ============================================================================
# SLIDE 19: Recommendations
# ============================================================================
print("  [19/20] Creating recommendations slide...")
slide = add_content_slide(prs, "Recommendations for Action")
bullets = [
    "For Policymakers:",
    "  - Prioritize Security & Rule of Law - foundation for other improvements",
    "  - Learn from improvers - study Gambia, Somalia, Seychelles reforms",
    "  - Monitor backsliding - early intervention in declining countries",
    "",
    "For Development Partners:",
    "  - Target Central Africa - greatest need, significant impact potential",
    "  - Support post-conflict recovery - Somalia proves ROI",
    "  - Strengthen regional institutions - integration correlates with performance",
    "",
    "For All Stakeholders:",
    "  - Take integrated approach - categories are interconnected",
    "  - Sustain long-term commitment - governance gains require persistence",
    "  - Share best practices - facilitate South-South learning"
]
add_bullet_points(slide, bullets, start_top=1.6)

# ============================================================================
# SLIDE 20: Closing Slide
# ============================================================================
print("  [20/20] Creating closing slide...")
slide = add_title_slide(prs, "Thank You", "Questions & Discussion")

# Add contact info
contact_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
contact_frame = contact_box.text_frame
contact_frame.text = "Data Source: Mo Ibrahim Foundation\nIbrahim Index of African Governance (IIAG) 2024\nhttps://iiag.online"
contact_para = contact_frame.paragraphs[0]
contact_para.alignment = PP_ALIGN.CENTER
contact_para.font.size = Pt(16)
contact_para.font.color.rgb = WHITE

# Save presentation
output_path = Path('IIAG_Presentation_2023.pptx')
prs.save(str(output_path))

print("\n" + "=" * 60)
print("PRESENTATION CREATED SUCCESSFULLY!")
print("=" * 60)
print(f"\nFile: {output_path.absolute()}")
print(f"Slides: 20")
print(f"Duration: 15-20 minutes")
print("\nSlide Breakdown:")
print("  1. Title slide")
print("  2. Agenda")
print("  3-4. IIAG overview & categories")
print("  5. Key statistics")
print("  6-8. Top performers analysis")
print("  9-11. Improvement champions")
print("  12-13. Governance challenges")
print("  14-15. Regional analysis")
print("  16-17. Category performance")
print("  18. Key insights")
print("  19. Recommendations")
print("  20. Closing/Q&A")
print("\nPresentation Features:")
print("  - Professional design with consistent color scheme")
print("  - Mix of data visualizations and bullet points")
print("  - Clear narrative flow from overview to recommendations")
print("  - Embedded high-resolution charts from analysis")
print("  - Ready for immediate presentation!")
